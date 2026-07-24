from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


COLORS = {
    "ink": "17212B",
    "muted": "5B6670",
    "line": "D6DCE2",
    "blue": "2878B5",
    "teal": "2A9D8F",
    "gold": "D99B2B",
    "red": "C65353",
    "light_blue": "EAF3FA",
    "light_teal": "E8F5F2",
    "light_gold": "FBF2DF",
    "light_red": "F9EAEA",
    "light_gray": "F4F6F8",
    "white": "FFFFFF",
}


def rgb(name: str) -> RGBColor:
    return RGBColor.from_string(COLORS[name])


def add_text(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str,
    *,
    size: float = 15,
    color: str = "ink",
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.01)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_shape(
    slide,
    kind,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: str = "white",
    edge: str = "ink",
    line_width: float = 1.2,
    dashed: bool = False,
):
    shape = slide.shapes.add_shape(
        kind,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(edge)
    shape.line.width = Pt(line_width)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str = "",
    *,
    fill: str = "white",
    edge: str = "ink",
    size: float = 15,
    bold: bool = False,
    dashed: bool = False,
):
    shape = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        width,
        height,
        fill=fill,
        edge=edge,
        dashed=dashed,
    )
    if text:
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.margin_left = frame.margin_right = Inches(0.04)
        frame.margin_top = frame.margin_bottom = Inches(0.02)
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb("ink")
    return shape


def add_arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = "muted",
    width: float = 1.2,
    dashed: bool = False,
    arrow: bool = True,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        line = connector.line._get_or_add_ln()
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        tail.set("w", "sm")
        tail.set("len", "sm")
        line.append(tail)
    return connector


def add_circle(slide, x: float, y: float, diameter: float, *, fill: str, edge: str):
    return add_shape(
        slide,
        MSO_SHAPE.OVAL,
        x,
        y,
        diameter,
        diameter,
        fill=fill,
        edge=edge,
        line_width=1.0,
    )


def add_panel(slide, x: float, title: str, letter: str, color: str) -> None:
    add_box(slide, x, 0.88, 4.0, 6.20, fill="white", edge="line")
    add_circle(slide, x + 0.22, 1.08, 0.36, fill=color, edge=color)
    add_text(slide, x + 0.22, 1.08, 0.36, 0.36, letter, size=15, color="white", bold=True)
    add_text(slide, x + 0.70, 1.06, 2.98, 0.40, title, size=17, bold=True, align=PP_ALIGN.LEFT)


def add_request_card(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    accent: str,
    repeated_marker: bool = True,
) -> None:
    add_box(slide, x, y, width, height, fill="white", edge=accent)
    for offset, fraction in ((0.22, 0.68), (0.39, 0.82), (0.56, 0.56)):
        line = add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            x + 0.16,
            y + offset,
            width * fraction,
            0.055,
            fill="line",
            edge="line",
            line_width=0.3,
        )
        line.rotation = 0
    if repeated_marker:
        add_shape(
            slide,
            MSO_SHAPE.DIAMOND,
            x + width - 0.38,
            y + 0.17,
            0.18,
            0.18,
            fill="light_gold",
            edge="gold",
            line_width=0.9,
        )


def add_agent(slide, x: float, y: float, color: str) -> None:
    add_circle(slide, x, y, 0.28, fill="white", edge=color)
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x - 0.10,
        y + 0.31,
        0.48,
        0.31,
        fill=f"light_{color}",
        edge=color,
        line_width=1.0,
    )


def add_id_tag(slide, x: float, y: float, label: str) -> None:
    add_box(slide, x, y, 0.82, 0.28, label, fill="light_gray", edge="line", size=11.5)
    add_arrow(slide, x + 0.08, y + 0.24, x + 0.74, y + 0.04, color="red", width=1.5, arrow=False)


def add_cluster(slide, x: float, y: float, color: str, marker: str) -> None:
    ring = add_shape(
        slide,
        MSO_SHAPE.OVAL,
        x,
        y,
        1.05,
        0.82,
        fill="white",
        edge=color,
        line_width=1.3,
    )
    ring.fill.transparency = 100000
    for dx, dy in ((0.18, 0.20), (0.46, 0.12), (0.67, 0.38)):
        if marker == "diamond":
            add_shape(slide, MSO_SHAPE.DIAMOND, x + dx, y + dy, 0.16, 0.16, fill=f"light_{color}", edge=color)
        else:
            add_circle(slide, x + dx, y + dy, 0.16, fill=f"light_{color}", edge=color)


def add_workflow(slide, x: float, y: float, color: str) -> None:
    for index in range(3):
        add_circle(slide, x + index * 0.37, y, 0.18, fill="white", edge=color)
        if index < 2:
            add_arrow(slide, x + 0.18 + index * 0.37, y + 0.09, x + 0.35 + index * 0.37, y + 0.09, color=color, width=0.9)


def build_presentation(output: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("white")

    add_text(slide, 0.35, 0.12, 12.63, 0.42, "From Unidentified Requests to Longitudinal Linkage", size=22, bold=True)
    add_text(slide, 0.35, 0.50, 12.63, 0.25, "Read left to right: observe what remains, reconstruct hidden groups, then follow them over time.", size=13.5, color="muted")

    # 1. Threat and controlled measurement.
    add_text(slide, 0.35, 0.84, 3.2, 0.28, "1  PROVIDER VIEW + PAIRED MEASUREMENT", size=14, color="muted", bold=True, align=PP_ALIGN.LEFT)
    for y, color in ((1.34, "blue"), (1.78, "teal"), (2.22, "gold")):
        add_agent(slide, 0.50, y, color)
        add_arrow(slide, 0.99, y + 0.29, 1.25, y + 0.29, color=color)
    add_box(slide, 1.28, 1.40, 1.20, 1.25, fill="light_blue", edge="blue")
    add_text(slide, 1.38, 1.50, 1.00, 0.28, "Broker", size=16, bold=True)
    add_id_tag(slide, 1.46, 2.08, "IDs")
    add_arrow(slide, 2.52, 2.02, 2.85, 2.02, color="blue", width=1.4)
    for offset in (0.0, 0.13, 0.26):
        add_request_card(slide, 2.86 + offset, 1.52 + offset, 1.02, 0.82, accent="teal")
    add_text(slide, 2.66, 2.43, 1.76, 0.34, "IDs disappear;\ncontent repeats", size=12.5, color="teal", bold=True)

    add_arrow(slide, 4.24, 2.02, 4.58, 2.02, color="muted", width=1.4)
    add_box(slide, 4.62, 1.30, 3.25, 1.48, fill="light_gray", edge="line")
    add_text(slide, 4.78, 1.37, 1.17, 0.25, "Original", size=13.5, color="gold", bold=True)
    add_request_card(slide, 4.88, 1.72, 0.78, 0.62, accent="gold")
    add_text(slide, 5.73, 1.84, 0.30, 0.25, "↔", size=18, color="gold", bold=True)
    add_request_card(slide, 6.08, 1.72, 0.78, 0.62, accent="muted", repeated_marker=False)
    add_text(slide, 6.03, 1.37, 1.60, 0.25, "Paired intervention", size=13.5, color="gold", bold=True)
    for index, label in enumerate(("handles", "replay", "timing / collision")):
        add_box(slide, 4.83 + index * 0.92, 2.42, 0.84, 0.25, label, fill="white", edge="gold", size=9.5)

    add_arrow(slide, 7.92, 2.02, 8.23, 2.02, color="muted", width=1.4)
    add_box(slide, 8.28, 1.39, 1.34, 1.18, "Run the same\nCARP / ASL", fill="light_red", edge="red", size=13.5, bold=True)
    add_arrow(slide, 9.66, 2.02, 10.00, 2.02, color="muted", width=1.4)
    add_box(slide, 10.05, 1.34, 2.53, 1.28, "Did linkage drop?\nIdentify the signal", fill="light_gray", edge="ink", size=14, bold=True)
    add_box(slide, 10.22, 2.64, 1.58, 0.41, "sealed truth: scoring only", fill="white", edge="red", size=9.5, dashed=True)
    add_arrow(slide, 11.02, 2.63, 11.30, 2.61, color="red", dashed=True)

    # 2. Complementary bounded methods.
    add_arrow(slide, 0.35, 3.20, 12.98, 3.20, color="line", width=0.8, arrow=False)
    add_text(slide, 0.35, 3.32, 3.4, 0.28, "2  TWO WAYS TO RECONSTRUCT HIDDEN GROUPS", size=14, color="muted", bold=True, align=PP_ALIGN.LEFT)
    for index, color in enumerate(("blue", "teal", "gold")):
        add_request_card(slide, 0.48 + index * 0.17, 3.85 + index * 0.13, 0.95, 0.74, accent=color)
    add_text(slide, 0.38, 4.70, 1.44, 0.25, "visible evidence", size=12.5, color="muted", bold=True)
    add_arrow(slide, 1.66, 4.33, 2.02, 4.33, color="muted", width=1.4)

    add_box(slide, 2.05, 3.76, 1.03, 0.62, "CARP", fill="light_blue", edge="blue", size=16, bold=True)
    for x, title, explanation in (
        (3.30, "Block & index", "find likely pairs"),
        (5.24, "Refine candidates", "verify continuity"),
        (7.18, "Propagate handles", "connect workflows"),
    ):
        add_box(slide, x, 3.70, 1.62, 0.75, fill="light_blue", edge="blue")
        add_text(slide, x + 0.05, 3.76, 1.52, 0.25, title, size=12, bold=True)
        add_text(slide, x + 0.05, 4.04, 1.52, 0.22, explanation, size=11.5, color="muted")
    for x1, x2 in ((3.09, 3.26), (4.94, 5.20), (6.88, 7.14)):
        add_arrow(slide, x1, 4.08, x2, 4.08, color="blue")

    add_box(slide, 2.05, 4.73, 1.03, 0.62, "ASL", fill="light_teal", edge="teal", size=16, bold=True)
    for x, title, explanation in (
        (3.30, "Agent state", "track bounded history"),
        (5.24, "Support vs conflict", "weigh evidence"),
        (7.18, "Selective hierarchy", "link or abstain"),
    ):
        add_box(slide, x, 4.67, 1.62, 0.75, fill="light_teal", edge="teal")
        add_text(slide, x + 0.05, 4.73, 1.52, 0.25, title, size=12, bold=True)
        add_text(slide, x + 0.05, 5.01, 1.52, 0.22, explanation, size=11.5, color="muted")
    for x1, x2 in ((3.09, 3.26), (4.94, 5.20), (6.88, 7.14)):
        add_arrow(slide, x1, 5.05, x2, 5.05, color="teal")

    add_arrow(slide, 8.84, 4.08, 9.20, 4.32, color="blue")
    add_arrow(slide, 8.84, 5.05, 9.20, 4.62, color="teal")
    add_cluster(slide, 9.22, 3.88, "blue", "circle")
    add_cluster(slide, 10.38, 4.52, "teal", "diamond")
    add_text(slide, 9.25, 5.33, 2.28, 0.26, "Recovered workflows and entities", size=13, bold=True)
    add_arrow(slide, 11.45, 4.55, 11.63, 4.55, color="muted", dashed=True)
    add_box(slide, 11.67, 3.86, 1.14, 1.36, "Avoid all-pairs\nreject weak links", fill="light_gray", edge="ink", size=12.5, bold=True)

    # 3. Channels and consequences.
    add_arrow(slide, 0.35, 5.72, 12.98, 5.72, color="line", width=0.8, arrow=False)
    add_text(slide, 0.35, 5.84, 3.0, 0.28, "3  WHAT EACH SIGNAL REVEALS", size=14, color="muted", bold=True, align=PP_ALIGN.LEFT)
    for x, label, width, edge, fill in (
        (0.48, "Direct exposure\nread fields", 1.54, "red", "light_red"),
        (2.17, "Continuity\ngroup one task", 1.72, "blue", "light_blue"),
        (4.04, "Propagation\nconnect later tasks", 2.00, "gold", "light_gold"),
    ):
        add_box(slide, x, 6.20, width, 0.60, label, fill=fill, edge=edge, size=12, bold=True)

    add_arrow(slide, 6.10, 6.50, 6.58, 6.50, color="muted", width=1.2)
    add_workflow(slide, 6.65, 6.40, "blue")
    add_arrow(slide, 7.65, 6.49, 8.04, 6.49, color="blue")
    add_shape(slide, MSO_SHAPE.HEXAGON, 8.08, 6.12, 0.76, 0.76, fill="light_gold", edge="gold")
    add_shape(slide, MSO_SHAPE.DIAMOND, 8.34, 6.38, 0.24, 0.24, fill="white", edge="gold")
    add_text(slide, 7.92, 6.91, 1.10, 0.24, "same project / customer", size=11.5, color="gold", bold=True)
    add_arrow(slide, 8.87, 6.34, 9.22, 6.10, color="gold")
    add_arrow(slide, 8.87, 6.65, 9.22, 6.88, color="gold")
    add_box(slide, 9.26, 5.96, 1.22, 0.55, "Aggregate\nPartial profile", fill="light_teal", edge="teal", size=12.5, bold=True)
    add_box(slide, 9.26, 6.72, 1.22, 0.55, "Recognize later\nWatchlist", fill="light_red", edge="red", size=12.5, bold=True)
    add_arrow(slide, 10.53, 6.49, 10.87, 6.49, color="muted")
    add_box(slide, 10.91, 6.08, 1.90, 0.84, "New requests join\nthe earlier component", fill="light_gold", edge="gold", size=13.5, bold=True)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate the editable visual main-paper figure.")
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("docs/overleaf/figures/provider_linkage_overview_editable.pptx"),
    )
    args = parser.parse_args()
    build_presentation(args.output)
    print(args.output)


if __name__ == "__main__":
    main()
